"""PowerPoint loader (python-pptx).

One ParsedPart per slide. A slide's text is assembled in reading order:
the title first, then every shape's text frame, then any table cells flattened
row-by-row, and finally the speaker notes (prefixed so they're distinguishable
from on-slide content). Empty slides yield an empty-text part — the embedding
step skips those, same as image-only PDF pages.

Structure detection (RAG_ROADMAP 3.1): each slide's title is emitted as a
level-1 heading at char offset 0, so the chunker treats slide boundaries as
hard section boundaries and builds a "Deck > Slide title" section path. Slides
without a title placeholder fall back to "Slide N" so the boundary still exists.

python-pptx reads the .pptx OOXML package directly (no COM / no PowerPoint
install needed), so this runs headless on the worker.
"""

from collections.abc import Iterator
from pathlib import Path

from pynote_core.parsers.types import ParsedPart

# Speaker notes are pedagogically useful (they're often the actual script) but
# they aren't visible on the slide. Prefix them so retrieval/citation can tell
# on-slide text from notes, and so the marker survives chunking.
_NOTES_PREFIX = "[Speaker notes] "


def _iter_shape_text(shapes) -> Iterator[str]:  # noqa: ANN001 - pptx shape collection
    """Yield text from each shape in a slide, flattening tables into rows.

    Shapes with a text frame contribute their text as-is. Table shapes are
    serialized cell-by-cell, tab-joined per row, so tabular data stays on one
    line and reads sensibly once embedded. Group shapes are walked recursively
    because PowerPoint nests placeholders inside groups more often than not.
    """
    for shape in shapes:
        # Groups nest other shapes — recurse before anything else.
        if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
            yield from _iter_shape_text(shape.shapes)
            continue
        if shape.has_table:
            for row in shape.table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                line = "\t".join(c for c in cells if c)
                if line:
                    yield line
            continue
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                yield text


def _slide_title(slide, ordinal: int) -> str:  # noqa: ANN001 - pptx slide
    """Best-effort slide title: the title placeholder, else a synthetic label.

    A stable non-empty title matters — it becomes the section heading, so even
    untitled slides need a deterministic label to anchor their section path.
    """
    title_shape = slide.shapes.title
    if title_shape is not None:
        title = (title_shape.text or "").strip()
        if title:
            return title
    return f"Slide {ordinal + 1}"


def _notes_text(slide) -> str:  # noqa: ANN001 - pptx slide
    """Speaker-notes text for a slide, or '' when the slide has no notes."""
    if not slide.has_notes_slide:
        return ""
    frame = slide.notes_slide.notes_text_frame
    if frame is None:
        return ""
    return (frame.text or "").strip()


def parse_pptx(path: Path) -> Iterator[ParsedPart]:
    """Yield one ParsedPart per slide, in presentation order.

    `page` is set to the 1-based slide number so citations can read "Slide 4"
    the same way PDF parts read "Page 4". `headings` carries the slide title as
    a single level-1 heading anchored at offset 0.
    """
    # Imported lazily so the dependency is only required on the worker actually
    # parsing a deck — the API and other loaders don't pull pptx in.
    from pptx import Presentation

    prs = Presentation(str(path))

    for ordinal, slide in enumerate(prs.slides):
        title = _slide_title(slide, ordinal)

        # Title leads so the heading offset is a real 0; body follows.
        blocks: list[str] = [title]
        blocks.extend(_iter_shape_text(slide.shapes))
        notes = _notes_text(slide)
        if notes:
            blocks.append(_NOTES_PREFIX + notes)

        # De-dupe the title if it also came back as a normal text-frame shape
        # (PowerPoint stores the title placeholder as a shape too).
        deduped: list[str] = []
        for b in blocks:
            if b not in deduped:
                deduped.append(b)
        text = "\n\n".join(deduped)

        yield ParsedPart(
            ordinal=ordinal,
            text=text,
            page=ordinal + 1,
            headings=[{"text": title, "level": 1, "start": 0}],
        )
